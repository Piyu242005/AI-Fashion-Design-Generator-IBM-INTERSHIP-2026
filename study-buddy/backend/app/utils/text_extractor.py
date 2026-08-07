"""
Text Extractor — AI-Powered Study Buddy
========================================
Extracts raw text from PDF, DOCX, PPTX, and TXT files.
Returns cleaned plain text ready for chunking.
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger("study_buddy.extractor")


def extract_text(file_path: str | Path) -> str:
    """
    Extract plain text from a file based on its extension.

    Args:
        file_path: Path to the uploaded file.

    Returns:
        Extracted text string.

    Raises:
        ValueError: If the file extension is not supported.
        RuntimeError: If extraction fails.
    """
    path = Path(file_path)
    ext  = path.suffix.lower().lstrip(".")

    extractors = {
        "pdf":  _extract_pdf,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "txt":  _extract_txt,
    }

    if ext not in extractors:
        raise ValueError(f"Unsupported file type: .{ext}")

    try:
        text = extractors[ext](path)
        cleaned = _clean_text(text)
        logger.info("Extracted %d chars from %s", len(cleaned), path.name)
        return cleaned
    except Exception as exc:
        logger.error("Extraction failed for %s: %s", path.name, exc)
        raise RuntimeError(f"Could not extract text from {path.name}: {exc}") from exc


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF (fitz)."""
    import fitz  # PyMuPDF

    text_parts: list[str] = []
    with fitz.open(str(path)) as doc:
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"[Page {page_num}]\n{page_text}")
    return "\n\n".join(text_parts)


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return "\n\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs   = Presentation(str(path))
    slides: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            slides.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    return "\n\n".join(slides)


def _extract_txt(path: Path) -> str:
    """Read a plain text file with UTF-8 encoding, fallback to latin-1."""
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


# ---------------------------------------------------------------------------
# Text cleaning
# ---------------------------------------------------------------------------

def _clean_text(text: str) -> str:
    """Remove excessive whitespace and non-printable characters."""
    import re

    # Normalise line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Remove non-printable chars (except newline, tab)
    text = re.sub(r"[^\x09\x0A\x20-\x7E\u00A0-\uFFFF]", " ", text)
    # Collapse multiple spaces
    text = re.sub(r" {2,}", " ", text)
    return text.strip()
